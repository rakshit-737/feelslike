"""Refresh docs/ from the measured result files — surgically, never by rebuilding.

The presentation in docs/FeelsLike_Presentation.pptx is a hand-designed 14-slide
deck (every element an absolutely positioned shape). Regenerating it from scratch
would throw away that design, so this script does three narrow things instead:

  * replace TEXT inside existing shapes while preserving the run's formatting;
  * clone an existing slide (deep-copying its shape XML) when a new slide is
    genuinely needed, so a new slide inherits the deck's visual language;
  * renumber the "NN / TOTAL" footers after any insertion.

Every number it writes comes from a result file on disk — evals/results_*.json and
rl/models/progress.csv — so the deck can never drift from what the code measured.
Capability claims are gated on CAPABILITIES below: a claim is only written once the
corresponding subsystem has actually been verified, which keeps the deck honest
while the build is still in flight.

Usage:
  python -m scripts.update_docs --check          # report what would change, write nothing
  python -m scripts.update_docs                  # update the deck + rebuild the report
  python -m scripts.update_docs --pptx-only
  python -m scripts.update_docs --report-only
"""
from __future__ import annotations

import argparse
import copy
import json
import re
import shutil
from pathlib import Path

from pptx.util import Inches

ROOT = Path(__file__).resolve().parent.parent
DOCS = ROOT / "docs"
DECK = DOCS / "FeelsLike_Presentation.pptx"


# ---------------------------------------------------------------------------
# capability gate — flip a flag to True only when the subsystem is VERIFIED
# ---------------------------------------------------------------------------
# Each entry is (verified, the sentence the deck should say). While verified is
# False the deck keeps its current, still-true wording. This is the mechanism
# that stops an in-flight feature from being announced as shipped.
CAPABILITIES: dict = {
    "multi_zone":     (False, "One complaint can name several zones; each gets its own constraint."),
    "explainability": (False, "Every setpoint change carries a structured, inspectable decision record."),
    "what_if":        (False, "Counterfactual scenarios run on a cloned twin, never on live state."),
    "maintenance":    (False, "Capacity, sensor, actuator and recurring faults are detected from evidence."),
    "humidity":       (False, "Zone humidity, dew point and latent load are modelled and reported."),
    "safety_modes":   (False, "Five operating modes from fully automatic to maintenance lockout."),
}


def capability_facts() -> dict:
    """The verified subset of CAPABILITIES. INPUT none. OUTPUT {key: sentence}."""
    return {k: text for k, (ok, text) in CAPABILITIES.items() if ok}


# ---------------------------------------------------------------------------
# measured facts
# ---------------------------------------------------------------------------

def load_facts() -> dict:
    """Read every number the docs quote straight from the result files.

    INPUT: none (reads evals/*.json, rl/models/progress.csv).
    OUTPUT: dict of measured facts, all floats/ints/strings.
    SIDE EFFECTS: none.
    ERROR STATES: raises FileNotFoundError if results_energy.json is missing —
      that file is the spine of every claim, so failing loudly is correct.
      Optional inputs (LLM benchmark, what-if) degrade to None.
    """
    energy = json.loads((ROOT / "evals/results_energy.json").read_text())["results"]

    def row(name):
        for k, v in energy.items():
            if k.startswith(name):
                return v
        raise KeyError(name)

    static, react = row("Static"), row("Reactive")
    rules, ppo = row("FeelsLike (constraint-aware)"), energy.get("FeelsLike (PPO agent)")

    facts = {
        "static_kwh": static["kwh"], "static_viol": static["viol_min"],
        "react_kwh": react["kwh"], "react_viol": react["viol_min"],
        "react_saved": react["saved_pct_vs_baseline"],
        "rules_kwh": rules["kwh"], "rules_viol": rules["viol_min"],
        "rules_saved": rules["saved_pct_vs_baseline"],
        "ppo_kwh": ppo["kwh"] if ppo else None,
        "ppo_viol": ppo["viol_min"] if ppo else None,
        "ppo_saved": ppo["saved_pct_vs_baseline"] if ppo else None,
        "saved_kwh": round(static["kwh"] - rules["kwh"], 1),
        "saved_rs": round(static["cost_rs"] - rules["cost_rs"]),
        "saved_co2": round(static["co2_kg"] - rules["co2_kg"]),
    }

    nlp_path = ROOT / "evals/results_nlp.json"
    if nlp_path.exists():
        splits = json.loads(nlp_path.read_text())["splits"]
        facts["nlp_splits"] = splits
        facts["nlp_cases"] = sum(s["n"] for s in splits.values())
        ho = {k: v for k, v in splits.items() if k.startswith("heldout")}
        if ho:
            n = sum(v["n"] for v in ho.values())
            facts["rules_heldout_triple_pct"] = round(100 * sum(v["triple"] for v in ho.values()) / n)

    llm_path = ROOT / "evals/results_nlp_llm.json"
    if llm_path.exists():
        splits = json.loads(llm_path.read_text())["splits"]
        ho = {k: v for k, v in splits.items() if k.startswith("heldout")}
        if ho:
            n = sum(v["n"] for v in ho.values())
            facts["llm_heldout_triple_pct"] = round(100 * sum(v["triple"] for v in ho.values()) / n)

    curve = ROOT / "rl/models/progress.csv"
    if curve.exists():
        import csv
        rows = list(csv.DictReader(curve.open()))
        pts = [float(r["rollout/ep_rew_mean"]) for r in rows if r.get("rollout/ep_rew_mean")]
        if pts:
            facts["reward_first"], facts["reward_last"] = round(pts[0], 1), round(pts[-1], 1)
            facts["train_steps"] = int(rows[-1]["time/total_timesteps"])

    # The blind probe is the number the docs are allowed to quote. A held-out
    # split that someone has debugged against is no longer held out; this one is
    # never used for tuning, so it is the honest generalization figure.
    blind = ROOT / "evals/results_blind.json"
    if blind.exists():
        b = json.loads(blind.read_text())
        facts["blind_n"] = b["n"]
        facts["blind_triple_pct"] = b["pct"]["triple"]
        facts["blind_zoneset_pct"] = b["pct"]["zoneset"]
        facts["blind_det_pct"] = b["pct"]["det"]
        facts["blind_issue_pct"] = b["pct"]["issue"]
    blind_llm = ROOT / "evals/results_blind_llm.json"
    if blind_llm.exists():
        bl = json.loads(blind_llm.read_text())["pct"]
        facts["blind_llm"] = {"Zone": bl["zoneset"], "Issue": bl["issue"],
                              "Detection": bl["det"], "Exact triple": bl["triple"]}

    wif = ROOT / "evals/results_whatif.json"
    facts["whatif"] = json.loads(wif.read_text()) if wif.exists() else None
    return facts


# ---------------------------------------------------------------------------
# pptx surgery
# ---------------------------------------------------------------------------

def set_text(shape, new_text: str) -> bool:
    """Replace a shape's text, keeping the first run's formatting.

    INPUT: a shape with a text frame; the replacement string.
    OUTPUT: True when the text changed.
    SIDE EFFECTS: mutates the shape in place.
    ERROR STATES: returns False for shapes with no text frame or no runs
      (an empty placeholder has no run whose formatting we could inherit).
    Technique: write into run 0 and blank the remaining runs, rather than
    assigning to text_frame.text, which would discard font styling.
    """
    if not shape.has_text_frame:
        return False
    paras = shape.text_frame.paragraphs
    if not paras or not paras[0].runs:
        return False
    if shape.text_frame.text.strip() == new_text.strip():
        return False
    paras[0].runs[0].text = new_text
    for r in paras[0].runs[1:]:
        r.text = ""
    for p in paras[1:]:
        for r in p.runs:
            r.text = ""
    return True


def find_shapes(slide, predicate):
    """Every text shape on a slide whose text satisfies predicate(text)."""
    return [sh for sh in slide.shapes
            if sh.has_text_frame and predicate(sh.text_frame.text.strip())]


def replace_on_slide(slide, matcher, new_text: str, log: list, label: str) -> bool:
    """Find one shape by matcher and rewrite it. Records the edit in log."""
    hits = find_shapes(slide, matcher)
    if not hits:
        log.append(f"  MISS  {label}: no shape matched")
        return False
    old = hits[0].text_frame.text.strip()
    if set_text(hits[0], new_text):
        log.append(f"  edit  {label}: {old[:48]!r} -> {new_text[:48]!r}")
        return True
    log.append(f"  same  {label}: already current")
    return False


def clone_slide(prs, index: int):
    """Append a deep copy of slide[index] and return the new slide.

    INPUT: presentation, index of the slide to use as a visual template.
    OUTPUT: the newly appended slide (same layout, same shapes, same design).
    SIDE EFFECTS: appends a slide to the presentation.
    ERROR STATES: IndexError for an out-of-range index. Note python-pptx has no
      public clone API; copying each shape's XML element is the supported
      workaround for decks whose shapes are absolutely positioned, as this one's are.
    """
    src = prs.slides[index]
    dst = prs.slides.add_slide(src.slide_layout)
    for shape in list(dst.shapes):        # strip layout placeholders
        shape._element.getparent().remove(shape._element)
    for shape in src.shapes:
        dst.shapes._spTree.append(copy.deepcopy(shape._element))
    return dst


def renumber_footers(prs, log: list) -> int:
    """Rewrite every 'NN / TOTAL' page footer to match the real slide count."""
    total = len(prs.slides)
    pattern = re.compile(r"^\s*(\d{1,2})\s*/\s*(\d{1,2})\s*$")
    fixed = 0
    for i, slide in enumerate(prs.slides, start=1):
        for sh in find_shapes(slide, lambda t: bool(pattern.match(t))):
            want = f"{i:02d} / {total}"
            if set_text(sh, want):
                fixed += 1
    if fixed:
        log.append(f"  edit  footers: renumbered {fixed} of {total} slides")
    return fixed



def derive_bar_scale(slide_index: int = 9):
    """Derive (baseline_emu, emu_per_percent) for the comparison chart.

    Read from the PRISTINE backup deck, never from the live file: by the time the
    bars are resized their value labels have already been rewritten, so deriving
    the scale in place would compare a new label against an old height and
    conclude nothing needs to move. Reading the backup makes the derivation
    independent of edit order.

    INPUT: index of the chart slide. OUTPUT: (baseline, emu_per_pct) or None when
    the backup is absent or its geometry is unrecognisable. SIDE EFFECTS: none.
    """
    backup = DOCS / (DECK.stem + ".backup.pptx")
    src = backup if backup.exists() else DECK
    try:
        from pptx import Presentation as _P
        slide = _P(str(src)).slides[slide_index]
    except Exception:
        return None
    bars = [sh for sh in slide.shapes
            if sh.left is not None and sh.top is not None and sh.height
            and sh.left > Inches(6) and sh.height > Inches(0.3)
            and sh.width and Inches(0.5) < sh.width < Inches(1.6)
            and not (sh.has_text_frame and sh.text_frame.text.strip())]
    labels = [sh for sh in slide.shapes
              if sh.has_text_frame and sh.left is not None and sh.left > Inches(6)
              and re.fullmatch(r"\d{1,3}", sh.text_frame.text.strip())]
    if len(bars) < 4 or not labels:
        return None
    baseline = max(b.top + b.height for b in bars)
    ratios = []
    for b in bars:
        lab = min(labels, key=lambda s: abs(s.left - b.left))
        try:
            pct = int(lab.text_frame.text.strip())
        except ValueError:
            continue
        if pct > 0:
            ratios.append(b.height / pct)
    if not ratios:
        return None
    ratios.sort()
    return baseline, ratios[len(ratios) // 2]      # median guards one odd bar


def resize_nlp_bars(slide, group_values: dict, labels: dict, log: list) -> int:
    """Redraw the comparison bars so height matches the number printed above it.

    INPUT: the chart slide, {group_label: (rules_pct, llm_pct)}, the axis-label
      shapes keyed by group name, the change log.
    OUTPUT: count of bars moved.
    SIDE EFFECTS: mutates bar rectangles' top/height and their value labels' top.
    ERROR STATES: returns 0 and logs a MISS when the bar rectangles cannot be
      identified, leaving the slide untouched rather than half-edited.

    Geometry is DERIVED from the deck, not hardcoded: every bar in this chart
    shares a common baseline (top + height is constant), and the value scale is
    linear from zero, so scale = height / percent for any existing bar. Value
    labels keep their measured offset above the bar they annotate. Updating the
    printed number without moving the bar would leave a chart that contradicts
    its own labels — worse than not updating it at all.
    """
    bars = [sh for sh in slide.shapes
            if sh.left is not None and sh.top is not None and sh.height
            and sh.left > Inches(6) and sh.height > Inches(0.3)
            and sh.width and Inches(0.5) < sh.width < Inches(1.6)
            and not (sh.has_text_frame and sh.text_frame.text.strip())]
    if len(bars) < 8:
        log.append("  MISS  chart bars: expected 8 bar rectangles, found "
                   f"{len(bars)} — geometry left untouched")
        return 0

    # scale + baseline are derived from the backup (see derive_bar_scale)
    value_labels = [sh for sh in slide.shapes
                    if sh.has_text_frame
                    and re.fullmatch(r"\d{1,3}", sh.text_frame.text.strip())
                    and sh.left is not None and sh.left > Inches(6)]
    if not value_labels:
        log.append("  MISS  chart bars: no value labels to anchor the scale")
        return 0
    label_offset = None
    moved = 0
    derived = derive_bar_scale()
    if derived is None:
        log.append("  MISS  chart bars: could not derive the value scale")
        return 0
    baseline, emu_per_pct = derived

    for group, (rules_pct, llm_pct) in group_values.items():
        lab = labels.get(group)
        if lab is None:
            continue
        pair = sorted((b for b in bars),
                      key=lambda b: abs(b.left - lab.left))[:2]
        if len(pair) < 2:
            continue
        pair.sort(key=lambda b: b.left)
        for bar, pct in zip(pair, (rules_pct, llm_pct)):
            if pct is None:
                continue
            # derive the scale once, from this bar's current height and label
            near_lbl = min(value_labels, key=lambda s: abs(s.left - bar.left))
            if label_offset is None:
                label_offset = bar.top - near_lbl.top
            new_h = int(round(emu_per_pct * pct))
            new_top = int(round(baseline - new_h))
            if abs(new_h - bar.height) < Inches(0.02):
                continue
            bar.height, bar.top = new_h, new_top
            near_lbl.top = int(round(new_top - (label_offset or Inches(0.47))))
            moved += 1
    if moved:
        log.append(f"  edit  chart geometry: {moved} bar(s) resized to match their labels")
    return moved


def update_nlp_chart(slides, facts: dict, log: list) -> bool:
    """Rewrite the parser-comparison bar chart from the blind-probe results.

    INPUT: the deck's slides, the measured facts, the change log.
    OUTPUT: True when the chart was updated.
    SIDE EFFECTS: mutates bar value labels, and relabels the third axis group
      from "Severity" to "Detection".
    ERROR STATES: returns False (logging a MISS) when the chart's axis labels are
      not found, i.e. the deck was redesigned and these anchors no longer hold.

    Geometry: the chart is four groups of two bars. Each group's axis label shares
    an x-origin with its RULES bar; the LLM bar sits ~1.1in to the right. So each
    value shape is matched by x-proximity to its group label, which is stable under
    edits in a way "nearest number" is not. Severity becomes Detection because the
    blind probe measures detection and does not score severity — relabelling is
    honest, leaving a stale severity bar is not.
    """
    if "blind_triple_pct" not in facts:
        return False
    rules = {"Zone": facts["blind_zoneset_pct"], "Issue": facts["blind_issue_pct"],
             "Detection": facts["blind_det_pct"], "Exact triple": facts["blind_triple_pct"]}
    llm = facts.get("blind_llm") or {}

    for slide in slides:
        labels = {}
        for name in ("Zone", "Issue", "Severity", "Detection", "Exact triple"):
            hit = find_shapes(slide, lambda t, N=name: t == N)
            if hit:
                labels[name] = hit[0]
        if "Zone" not in labels or "Exact triple" not in labels:
            continue

        if "Severity" in labels and "Detection" not in labels:
            if set_text(labels["Severity"], "Detection"):
                log.append("  edit  chart axis: 'Severity' -> 'Detection' (what the probe measures)")
            labels["Detection"] = labels.pop("Severity")

        values = [s for s in slide.shapes
                  if s.has_text_frame and re.fullmatch(r"\d{1,3}", s.text_frame.text.strip())
                  and (s.left or 0) > Inches(6)]
        for group, lab in labels.items():
            if group not in rules:
                continue
            near = sorted(values, key=lambda s: abs((s.left or 0) - (lab.left or 0)))[:2]
            if len(near) < 2:
                continue
            bar_rules, bar_llm = sorted(near, key=lambda s: (s.left or 0))
            if set_text(bar_rules, str(rules[group])):
                log.append(f"  edit  chart {group} (rules) -> {rules[group]}%")
            if group in llm and set_text(bar_llm, str(llm[group])):
                log.append(f"  edit  chart {group} (LLM) -> {llm[group]}%")

        resize_nlp_bars(slide, {g: (rules.get(g), llm.get(g)) for g in rules},
                        labels, log)

        caption = find_shapes(slide, lambda t: t.startswith("dev accuracy"))
        if caption:
            n = facts.get("blind_n", 20)
            set_text(caption[0], f"measured on {n} unseen cases — never used for tuning")
            log.append("  edit  chart caption -> blind-probe framing")
        return True

    log.append("  MISS  nlp chart: axis labels not found (deck redesigned?)")
    return False


def update_deck(facts: dict, caps: dict, dry_run: bool) -> list:
    """Apply every data-driven edit the deck needs. Returns the change log."""
    from pptx import Presentation

    log: list = []
    prs = Presentation(str(DECK))
    slides = list(prs.slides)

    def match_format(old: str, value: float) -> str:
        """Render value the way the deck already renders that field.

        The deck shows headline kWh as whole numbers; rewriting 722 as 722.4
        would be pure churn that also breaks the visual rhythm of the bar labels.
        So the incumbent formatting wins unless the value genuinely changed.
        """
        return f"{value:.0f}" if "." not in old else f"{value:.1f}"

    # --- energy comparison bars (slide 8 in the current deck) --------------
    def edit_once(matcher, build, label):
        """Apply one logical edit across the deck; log a single MISS if absent."""
        for slide in slides:
            for sh in find_shapes(slide, matcher):
                old = sh.text_frame.text.strip()
                new = build(old)
                if new is None:
                    continue
                if set_text(sh, new):
                    log.append(f"  edit  {label}: {old[:44]!r} -> {new[:44]!r}")
                return True
        log.append(f"  MISS  {label}: no shape matched (deck wording may have changed)")
        return False

    # --- energy comparison bars: the value shape nearest each bar label ----
    bar_values = {
        "Static 22": facts["static_kwh"],
        "Reactive 24": facts["react_kwh"],
        "FeelsLike \u00b7 rules": facts["rules_kwh"],
        "FeelsLike \u00b7 PPO": facts["ppo_kwh"],
    }
    for label, value in bar_values.items():
        if value is None:
            continue
        for slide in slides:
            labels = find_shapes(slide, lambda t, L=label: t.startswith(L))
            if not labels:
                continue
            lab = labels[0]
            numeric = [sh for sh in slide.shapes
                       if sh.has_text_frame
                       and re.fullmatch(r"\d{2,4}(\.\d)?", sh.text_frame.text.strip())]
            if not numeric:
                continue
            near = min(numeric, key=lambda sh: abs((sh.top or 0) - (lab.top or 0)))
            old = near.text_frame.text.strip()
            if set_text(near, match_format(old, value)):
                log.append(f"  edit  bar {label}: {old} -> {match_format(old, value)} kWh")
            break

    # --- NLP: the deck must quote the BLIND number, never the tuned split --
    # The matcher deliberately anchors on "Exact triple" alone: a previous run
    # already rewrote this shape away from its original "Exact triple, held-out"
    # wording, and a matcher tied to the original text would then silently stop
    # matching — leaving a stale percentage on the slide with only a MISS in the log.
    if "blind_triple_pct" in facts:
        edit_once(lambda t: t.startswith("Exact triple"),
                  lambda _o: (f"Exact triple on unseen phrasings: "
                              f"{facts['blind_triple_pct']}% (tuned split reads higher)."),
                  "nlp held-out sentence")
    elif "llm_heldout_triple_pct" in facts and "rules_heldout_triple_pct" in facts:
        edit_once(lambda t: t.startswith("Exact triple"),
                  lambda _o: (f"Exact triple, held-out: LLM {facts['llm_heldout_triple_pct']}% "
                              f"vs rules {facts['rules_heldout_triple_pct']}%."),
                  "nlp held-out sentence")

    # --- NLP slide furniture: split composition, headline, blind framing ----
    if "nlp_splits" in facts:
        splits = facts["nlp_splits"]
        dev_n = splits.get("dev", {}).get("n", 0)
        ho_n = sum(v["n"] for k, v in splits.items() if k.startswith("heldout"))
        blind_n = facts.get("blind_n")
        parts = [f"{dev_n} dev", f"{ho_n} held-out"]
        if blind_n:
            parts.append(f"{blind_n} blind probe, never tuned on")
        edit_once(lambda t: re.match(r"^\d+ dev\b", t) is not None,
                  lambda _o: " · ".join(parts),
                  "nlp split composition")

    if "blind_triple_pct" in facts:
        edit_once(lambda t: t.startswith("The held-out gap is"),
                  lambda _o: "The blind probe is the honest number",
                  "nlp slide headline")

    # --- limitations: the parser ceiling is measured, not a remembered range -
    if "blind_triple_pct" in facts:
        edit_once(lambda t: t.startswith("Held-out exact-triple NLP accuracy")
                  or t.startswith("Blind-probe exact-triple NLP accuracy"),
                  lambda _o: (f"Blind-probe exact-triple NLP accuracy is "
                              f"{facts['blind_triple_pct']}% on unseen phrasing."),
                  "limitation: parser ceiling")

    if "nlp_cases" in facts:
        # Anchor on GEOMETRY, not on "nearest number". An earlier version picked the
        # numeric shape closest in `top` to the "labeled cases" caption and hit a bar
        # value in the chart below it instead, silently corrupting the chart. The
        # headline stat sits in the slide's left column; bars live past x = 6in.
        for slide in slides:
            caption = find_shapes(slide, lambda t: t == "labeled cases")
            if not caption:
                continue
            cap = caption[0]
            candidates = [s for s in slide.shapes
                          if s.has_text_frame
                          and re.fullmatch(r"\d{2,3}", s.text_frame.text.strip())
                          and (s.left or 0) < Inches(6)
                          and abs((s.left or 0) - (cap.left or 0)) < Inches(1.5)]
            if not candidates:
                log.append("  MISS  benchmark case count: no stat shape in the left column")
                continue
            near = min(candidates, key=lambda s: abs((s.top or 0) - (cap.top or 0)))
            if set_text(near, str(facts["nlp_cases"])):
                log.append(f"  edit  benchmark case count -> {facts['nlp_cases']}")

    update_nlp_chart(slides, facts, log)

    # --- RL reward sentence ----------------------------------------------
    if "reward_first" in facts:
        steps_m = facts.get("train_steps", 0) / 1_000_000
        edit_once(lambda t: t.startswith("Mean reward"),
                  lambda _o: (f"Mean reward \u2212{abs(facts['reward_first'])} \u2192 "
                              f"\u2212{abs(facts['reward_last'])}, still climbing at "
                              f"{steps_m:.0f}M steps."),
                  "rl reward sentence")

    # --- scale numbers ----------------------------------------------------
    edit_once(lambda t: re.fullmatch(r"[\d,]+ kWh", t) is not None,
              lambda _o: f"{facts['saved_kwh']:.0f} kWh", "measured kWh saved")
    edit_once(lambda t: t.replace(" ", "").startswith("\u20b91,"),
              lambda _o: f"\u20b9{facts['saved_rs']:,}", "measured cost saved")

    # --- limitations: retire a limitation once its fix is VERIFIED --------
    retired = {
        "multi_zone": "Multi-zone complaints currently carry a single zone.",
        "humidity": "Humidity is not modelled.",
    }
    for key, old_line in retired.items():
        if key not in caps:
            continue
        for slide in slides:
            for sh in find_shapes(slide, lambda t, o=old_line: t.strip() == o):
                if set_text(sh, caps[key]):
                    log.append(f"  edit  limitation retired ({key})")

    renumber_footers(prs, log)

    if not log:
        log.append("  deck already matches every measured value — nothing to change")
    if not dry_run:
        target = DECK
        backup = DOCS / (DECK.stem + ".backup.pptx")
        shutil.copy2(DECK, backup)
        try:
            prs.save(str(target))
            log.append(f"  saved  {target.name} (backup: {backup.name})")
        except PermissionError:
            alt = DOCS / (DECK.stem + ".updated.pptx")
            prs.save(str(alt))
            log.append(f"  LOCKED {target.name} is open in PowerPoint — wrote {alt.name} instead")
    return log


# ---------------------------------------------------------------------------
# report
# ---------------------------------------------------------------------------

def update_report(dry_run: bool) -> list:
    """Rebuild the HTML report (and PDF when a headless browser is available)."""
    log = []
    if dry_run:
        return ["  would run scripts.build_report and re-print the PDF"]
    from scripts import build_report
    build_report.main()
    log.append("  rebuilt docs/FeelsLike-Report.html from the result files")

    edge = Path("C:/Program Files (x86)/Microsoft/Edge/Application/msedge.exe")
    if not edge.exists():
        edge = Path("C:/Program Files/Microsoft/Edge/Application/msedge.exe")
    if edge.exists():
        import subprocess
        out = DOCS / "FeelsLike-Report.pdf"
        subprocess.run([str(edge), "--headless", "--disable-gpu", "--no-pdf-header-footer",
                        f"--print-to-pdf={out}", (DOCS / 'FeelsLike-Report.html').as_uri()],
                       capture_output=True, timeout=180)
        log.append(f"  printed {out.name} ({out.stat().st_size // 1024} KB)")
    else:
        log.append("  no headless browser found — HTML updated, PDF not re-printed")
    return log


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--check", action="store_true", help="report changes, write nothing")
    ap.add_argument("--pptx-only", action="store_true")
    ap.add_argument("--report-only", action="store_true")
    args = ap.parse_args()

    facts = load_facts()
    caps = capability_facts()

    print("MEASURED FACTS")
    for k, v in facts.items():
        if k not in ("nlp_splits", "whatif"):
            print(f"  {k:26} {v}")
    print(f"\nVERIFIED CAPABILITY CLAIMS: {len(caps)}/{len(CAPABILITIES)}")
    for k in CAPABILITIES:
        print(f"  [{'x' if k in caps else ' '}] {k}")
    pending = [k for k in CAPABILITIES if k not in caps]
    if pending:
        print(f"  -> still gated (not written into the docs): {', '.join(pending)}")

    if not args.report_only:
        print("\nDECK")
        for line in update_deck(facts, caps, args.check):
            print(line)
    if not args.pptx_only:
        print("\nREPORT")
        for line in update_report(args.check):
            print(line)


if __name__ == "__main__":
    main()
